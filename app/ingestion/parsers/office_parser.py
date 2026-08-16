"""Office & text document parser adapter (DOCX, XLSX, PPTX, TXT, MD)."""

import time
from pathlib import Path

import docx
import openpyxl
import pptx

from app.ingestion.parsers.base import (
    ElementType,
    ParsedDocument,
    ParsedElement,
    ParsedPage,
    ParsedTable,
)


class OfficeParser:
    """Parser for DOCX, XLSX, PPTX, and text formats."""

    parser_name: str = "office_parser"

    def parse(self, file_path: Path | str, mime_type: str | None = None) -> ParsedDocument:
        path = Path(file_path)
        ext = path.suffix.lower()

        if ext == ".docx":
            return self._parse_docx(path)
        elif ext in (".xlsx", ".xls"):
            return self._parse_xlsx(path)
        elif ext in (".pptx", ".ppt"):
            return self._parse_pptx(path)
        else:
            return self._parse_text(path)

    def _parse_docx(self, path: Path) -> ParsedDocument:
        start_time = time.perf_counter()
        doc = docx.Document(str(path))
        elements: list[ParsedElement] = []
        tables: list[ParsedTable] = []
        seq = 0

        for _p_idx, p in enumerate(doc.paragraphs):
            text = p.text.strip()
            if not text:
                continue
            seq += 1
            style_name = p.style.name.lower() if p.style else ""
            el_type = ElementType.PARAGRAPH
            level = None

            if "heading 1" in style_name or "title" in style_name:
                el_type = ElementType.HEADING
                level = 1
            elif "heading 2" in style_name:
                el_type = ElementType.HEADING
                level = 2
            elif "heading 3" in style_name:
                el_type = ElementType.HEADING
                level = 3
            elif "list" in style_name or "bullet" in style_name:
                el_type = ElementType.LIST

            elements.append(
                ParsedElement(
                    element_id=f"docx_elem_{seq}",
                    element_type=el_type,
                    text=text,
                    page_number=1,
                    level=level,
                    sequence_index=seq,
                )
            )

        for t_idx, table in enumerate(doc.tables):
            if not table.rows:
                continue
            headers = [c.text.strip() for c in table.rows[0].cells]
            data_rows = [[c.text.strip() for c in row.cells] for row in table.rows[1:]]

            md_lines = [
                "| " + " | ".join(headers) + " |",
                "| " + " | ".join(["---"] * len(headers)) + " |",
            ]
            for r in data_rows:
                md_lines.append("| " + " | ".join(r) + " |")

            tables.append(
                ParsedTable(
                    table_id=f"docx_table_{t_idx + 1}",
                    title=f"Table {t_idx + 1}",
                    page_number=1,
                    num_rows=len(data_rows),
                    num_cols=len(headers),
                    headers=headers,
                    cells=data_rows,
                    markdown="\n".join(md_lines),
                )
            )

        duration_ms = (time.perf_counter() - start_time) * 1000.0
        return ParsedDocument(
            filename=path.name,
            file_type="docx",
            total_pages=1,
            pages=[ParsedPage(page_number=1, elements=elements, tables=tables)],
            metadata={"source_path": str(path)},
            parser_name=self.parser_name,
            parsing_duration_ms=duration_ms,
        )

    def _parse_xlsx(self, path: Path) -> ParsedDocument:
        start_time = time.perf_counter()
        wb = openpyxl.load_workbook(str(path), data_only=True)
        pages: list[ParsedPage] = []

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            all_rows = list(ws.iter_rows(values_only=True))
            if not all_rows:
                continue

            headers = [str(c or "").strip() for c in all_rows[0]]
            data_rows = [
                [str(c if c is not None else "").strip() for c in row] for row in all_rows[1:]
            ]

            md_lines = [
                "| " + " | ".join(headers) + " |",
                "| " + " | ".join(["---"] * len(headers)) + " |",
            ]
            for r in data_rows:
                md_lines.append("| " + " | ".join(r) + " |")

            tables = [
                ParsedTable(
                    table_id=f"sheet_{sheet_idx + 1}_table",
                    title=f"Sheet: {sheet_name}",
                    page_number=sheet_idx + 1,
                    num_rows=len(data_rows),
                    num_cols=len(headers),
                    headers=headers,
                    cells=data_rows,
                    markdown="\n".join(md_lines),
                )
            ]

            elements = [
                ParsedElement(
                    element_id=f"sheet_header_{sheet_idx + 1}",
                    element_type=ElementType.HEADING,
                    text=f"Spreadsheet Sheet: {sheet_name}",
                    page_number=sheet_idx + 1,
                    level=1,
                    sequence_index=1,
                )
            ]

            pages.append(ParsedPage(page_number=sheet_idx + 1, elements=elements, tables=tables))

        duration_ms = (time.perf_counter() - start_time) * 1000.0
        return ParsedDocument(
            filename=path.name,
            file_type="xlsx",
            total_pages=len(pages),
            pages=pages,
            metadata={"source_path": str(path), "sheets": wb.sheetnames},
            parser_name=self.parser_name,
            parsing_duration_ms=duration_ms,
        )

    def _parse_pptx(self, path: Path) -> ParsedDocument:
        start_time = time.perf_counter()
        prs = pptx.Presentation(str(path))
        pages: list[ParsedPage] = []
        global_seq = 0

        for slide_idx, slide in enumerate(prs.slides):
            page_num = slide_idx + 1
            elements: list[ParsedElement] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        text = p.text.strip()
                        if text:
                            global_seq += 1
                            el_type = (
                                ElementType.HEADING
                                if shape == slide.shapes.title
                                else ElementType.PARAGRAPH
                            )
                            level = 1 if shape == slide.shapes.title else None
                            elements.append(
                                ParsedElement(
                                    element_id=f"slide_{page_num}_elem_{global_seq}",
                                    element_type=el_type,
                                    text=text,
                                    page_number=page_num,
                                    level=level,
                                    sequence_index=global_seq,
                                )
                            )

            pages.append(ParsedPage(page_number=page_num, elements=elements))

        duration_ms = (time.perf_counter() - start_time) * 1000.0
        return ParsedDocument(
            filename=path.name,
            file_type="pptx",
            total_pages=len(pages),
            pages=pages,
            metadata={"source_path": str(path)},
            parser_name=self.parser_name,
            parsing_duration_ms=duration_ms,
        )

    def _parse_text(self, path: Path) -> ParsedDocument:
        start_time = time.perf_counter()
        content = path.read_text(encoding="utf-8")
        lines = content.splitlines()
        elements: list[ParsedElement] = []

        for idx, line in enumerate(lines):
            clean = line.strip()
            if not clean:
                continue
            el_type = ElementType.PARAGRAPH
            level = None
            if clean.startswith("# "):
                el_type = ElementType.HEADING
                level = 1
            elif clean.startswith("## "):
                el_type = ElementType.HEADING
                level = 2
            elif clean.startswith("### "):
                el_type = ElementType.HEADING
                level = 3

            elements.append(
                ParsedElement(
                    element_id=f"text_elem_{idx + 1}",
                    element_type=el_type,
                    text=clean,
                    page_number=1,
                    level=level,
                    sequence_index=idx + 1,
                )
            )

        duration_ms = (time.perf_counter() - start_time) * 1000.0
        return ParsedDocument(
            filename=path.name,
            file_type="txt",
            total_pages=1,
            pages=[ParsedPage(page_number=1, elements=elements, raw_text=content)],
            metadata={"source_path": str(path)},
            parser_name=self.parser_name,
            parsing_duration_ms=duration_ms,
        )
