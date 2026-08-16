"""Unit tests for document parsers (Task 1.3, Task 1.4)."""

from pathlib import Path
import docx
import openpyxl
import pptx
import pytest

from app.ingestion.parsers.base import ElementType, ParsedDocument
from app.ingestion.parsers.docling_parser import DoclingParser
from app.ingestion.parsers.office_parser import OfficeParser
from app.ingestion.parsers.opendataloader_parser import OpenDataLoaderParser
from app.ingestion.parsers.pymupdf_parser import PyMuPDFParser

CORPUS_DIR = Path("benchmarks/corpus")


@pytest.fixture
def staff_handbook_pdf() -> Path:
    return CORPUS_DIR / "staff_handbook.pdf"


@pytest.fixture
def health_plan_pdf() -> Path:
    return CORPUS_DIR / "health_plan_at_a_glance_2026.pdf"


@pytest.fixture
def temp_docx(tmp_path: Path) -> Path:
    p = tmp_path / "sample_policy.docx"
    doc = docx.Document()
    doc.add_heading("Global Travel Policy", level=1)
    doc.add_paragraph("Employees are reimbursed for reasonable business travel.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Tier"
    table.rows[0].cells[1].text = "Allowance"
    r = table.add_row()
    r.cells[0].text = "Tier 1"
    r.cells[1].text = "$100"
    doc.save(str(p))
    return p


@pytest.fixture
def temp_xlsx(tmp_path: Path) -> Path:
    p = tmp_path / "bonus_matrix.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Bonus"
    ws.append(["Grade", "Target Bonus"])
    ws.append(["Grade 1", 0.05])
    wb.save(str(p))
    return p


@pytest.fixture
def temp_pptx(tmp_path: Path) -> Path:
    p = tmp_path / "presentation.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "DEI Strategy"
    slide.placeholders[1].text = "2026 Roadmap"
    prs.save(str(p))
    return p


def test_docling_parser_pdf(health_plan_pdf: Path) -> None:
    if not health_plan_pdf.exists():
        pytest.skip("Corpus file not found.")

    parser = DoclingParser()
    doc = parser.parse(health_plan_pdf)

    assert isinstance(doc, ParsedDocument)
    assert doc.parser_name == "docling"
    assert doc.total_pages == 11
    assert len(doc.pages) == 11
    assert len(doc.all_elements) > 0

    # Check bounding box presence
    for el in doc.all_elements[:20]:
        assert el.bounding_box is not None
        assert el.bounding_box.page_number >= 1


def test_opendataloader_parser_pdf(staff_handbook_pdf: Path) -> None:
    if not staff_handbook_pdf.exists():
        pytest.skip("Corpus file not found.")

    parser = OpenDataLoaderParser()
    doc = parser.parse(staff_handbook_pdf)

    assert isinstance(doc, ParsedDocument)
    assert doc.parser_name == "opendataloader"
    assert doc.total_pages == 56
    assert len(doc.all_elements) > 0


def test_pymupdf_parser_pdf(staff_handbook_pdf: Path) -> None:
    if not staff_handbook_pdf.exists():
        pytest.skip("Corpus file not found.")

    parser = PyMuPDFParser()
    doc = parser.parse(staff_handbook_pdf)

    assert isinstance(doc, ParsedDocument)
    assert doc.parser_name == "pymupdf"
    assert doc.total_pages == 56
    assert len(doc.all_elements) > 0


def test_office_parser_docx(temp_docx: Path) -> None:
    parser = OfficeParser()
    doc = parser.parse(temp_docx)

    assert isinstance(doc, ParsedDocument)
    assert doc.file_type == "docx"
    assert len(doc.all_elements) > 0
    assert len(doc.all_tables) == 1
    assert "Tier" in doc.all_tables[0].headers


def test_office_parser_xlsx(temp_xlsx: Path) -> None:
    parser = OfficeParser()
    doc = parser.parse(temp_xlsx)

    assert isinstance(doc, ParsedDocument)
    assert doc.file_type == "xlsx"
    assert len(doc.all_tables) >= 1
    assert "Grade" in doc.all_tables[0].headers


def test_office_parser_pptx(temp_pptx: Path) -> None:
    parser = OfficeParser()
    doc = parser.parse(temp_pptx)

    assert isinstance(doc, ParsedDocument)
    assert doc.file_type == "pptx"
    assert doc.total_pages == 1
    assert len(doc.all_elements) > 0
