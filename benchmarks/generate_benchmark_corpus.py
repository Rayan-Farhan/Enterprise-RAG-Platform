"""Generate representative HR benchmark corpus across multiple formats and stress dimensions."""

from pathlib import Path

import docx
import fitz  # PyMuPDF
import openpyxl
import pptx
from docx.shared import RGBColor
from openpyxl.styles import Alignment, Font, PatternFill
from PIL import Image, ImageDraw

CORPUS_DIR = Path(r"d:\Projects & Certificates\Projects\Enterprise-RAG-Platform\benchmarks\corpus")


def generate_comprehensive_hr_handbook(output_path: Path) -> None:
    """Generate a multi-page comprehensive enterprise HR policy handbook PDF."""
    doc = fitz.open()

    # Page 1: Cover & Table of Contents
    page1 = doc.new_page(width=595, height=842)
    page1.insert_text(
        (50, 80),
        "ENTERPRISE GLOBAL HR POLICY HANDBOOK",
        fontsize=18,
        fontname="helv",
        color=(0.1, 0.2, 0.5),
    )
    page1.insert_text(
        (50, 110),
        "Version: 2026.1 | Effective Date: January 1, 2026 | Authority: Global People Operations",
        fontsize=10,
        fontname="helv",
        color=(0.3, 0.3, 0.3),
    )

    page1.insert_text(
        (50, 160),
        "1. Introduction & Organizational Principles",
        fontsize=14,
        fontname="helv",
        color=(0.1, 0.2, 0.5),
    )
    page1.insert_text(
        (50, 185),
        "This Handbook sets forth the official employment standards, policies, and operational guidelines",
        fontsize=10,
        fontname="helv",
    )
    page1.insert_text(
        (50, 200),
        "applicable to all full-time, part-time, and contracted personnel across global operations.",
        fontsize=10,
        fontname="helv",
    )

    page1.insert_text(
        (50, 240),
        "2. Employment Categories & Probationary Periods",
        fontsize=14,
        fontname="helv",
        color=(0.1, 0.2, 0.5),
    )
    page1.insert_text(
        (50, 265),
        "2.1 Standard Probation Period",
        fontsize=12,
        fontname="helv",
        color=(0.2, 0.2, 0.2),
    )
    page1.insert_text(
        (50, 285),
        "All new employees join on an initial probationary status of exactly ninety (90) calendar days.",
        fontsize=10,
        fontname="helv",
    )
    page1.insert_text(
        (50, 300),
        "During probation, performance is formally reviewed at day 30, day 60, and day 85.",
        fontsize=10,
        fontname="helv",
    )
    page1.insert_text(
        (50, 315),
        "Probation may be extended up to a maximum of thirty (30) additional days upon written HR approval.",
        fontsize=10,
        fontname="helv",
    )

    # Table on Page 1
    page1.insert_text(
        (50, 350),
        "Table 1.1: Employment Category Classification",
        fontsize=11,
        fontname="helv",
        color=(0.1, 0.2, 0.5),
    )
    rect_table = fitz.Rect(50, 365, 545, 465)
    page1.draw_rect(rect_table, color=(0.8, 0.8, 0.8), fill=(0.95, 0.95, 0.98))

    headers = ["Category", "Standard Hours", "Overtime Eligible", "Notice Period"]
    x_positions = [60, 160, 280, 420]
    for x, h in zip(x_positions, headers, strict=False):
        page1.insert_text((x, 385), h, fontsize=10, fontname="helv", color=(0.1, 0.2, 0.4))

    rows = [
        ["Grade A (Executive)", "40 hrs/wk", "No (Exempt)", "90 Days"],
        ["Grade B (Professional)", "40 hrs/wk", "No (Exempt)", "60 Days"],
        ["Grade C (Technical/Ops)", "40 hrs/wk", "Yes (Non-Exempt)", "30 Days"],
        ["Grade D (Contractual)", "Variable", "No", "14 Days"],
    ]
    y = 405
    for row in rows:
        page1.draw_line(fitz.Point(50, y - 5), fitz.Point(545, y - 5), color=(0.85, 0.85, 0.85))
        for x, val in zip(x_positions, row, strict=False):
            page1.insert_text((x, y + 10), val, fontsize=9, fontname="helv")
        y += 18

    # Page 2: Leave Policy & Calculations
    page2 = doc.new_page(width=595, height=842)
    page2.insert_text(
        (50, 80),
        "3. Comprehensive Leave Entitlements",
        fontsize=14,
        fontname="helv",
        color=(0.1, 0.2, 0.5),
    )
    page2.insert_text(
        (50, 105),
        "3.1 Annual Paid Time Off (PTO)",
        fontsize=12,
        fontname="helv",
        color=(0.2, 0.2, 0.2),
    )
    page2.insert_text(
        (50, 125),
        "Employees accrue PTO monthly at a baseline rate of 1.66 days per month (20 days per annum).",
        fontsize=10,
        fontname="helv",
    )
    page2.insert_text(
        (50, 140),
        "Employees with five (5) or more years of tenure accrue at 2.08 days per month (25 days per annum).",
        fontsize=10,
        fontname="helv",
    )
    page2.insert_text(
        (50, 155),
        "A maximum of five (5) unused annual leave days may be carried forward into the next calendar year.",
        fontsize=10,
        fontname="helv",
    )
    page2.insert_text(
        (50, 170),
        "Carried forward leave must be utilized by March 31, or it is automatically forfeited without cash-out.",
        fontsize=10,
        fontname="helv",
    )

    page2.insert_text(
        (50, 205),
        "3.2 Parental & Maternity Leave",
        fontsize=12,
        fontname="helv",
        color=(0.2, 0.2, 0.2),
    )
    page2.insert_text(
        (50, 225),
        "Female employees are entitled to twenty-four (24) consecutive weeks of fully paid maternity leave.",
        fontsize=10,
        fontname="helv",
    )
    page2.insert_text(
        (50, 240),
        "Paternity and secondary caregiver leave provides six (6) consecutive weeks of fully paid leave.",
        fontsize=10,
        fontname="helv",
    )

    page2.insert_text(
        (50, 275),
        "3.3 Severance Calculation Formula",
        fontsize=12,
        fontname="helv",
        color=(0.2, 0.2, 0.2),
    )
    page2.insert_text(
        (50, 295),
        "In the event of redundancy, severance pay is calculated using the following statutory formula:",
        fontsize=10,
        fontname="helv",
    )

    # Formula representation
    page2.insert_text(
        (80, 320),
        "Severance Pay = (Years of Service * 1.5 * Monthly Base Salary) + Accrued PTO Balance",
        fontsize=10,
        fontname="helv",
        color=(0.6, 0.1, 0.1),
    )
    page2.insert_text(
        (50, 345),
        "Where partial years exceeding six (6) months round up to one full continuous year.",
        fontsize=9,
        fontname="helv",
        color=(0.3, 0.3, 0.3),
    )

    # Save PDF
    doc.save(str(output_path))
    doc.close()


def generate_table_heavy_compensation(output_path: Path) -> None:
    """Generate a table-heavy compensation and benefits grid PDF."""
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)

    page.insert_text(
        (50, 60),
        "GLOBAL COMPENSATION & MEDICAL BENEFIT TIERS 2026",
        fontsize=16,
        fontname="helv",
        color=(0.1, 0.2, 0.5),
    )
    page.insert_text(
        (50, 85),
        "Confidential — Human Resources Internal Distribution Only",
        fontsize=9,
        fontname="helv",
        color=(0.5, 0.1, 0.1),
    )

    # Table 1: Medical Plans
    page.insert_text(
        (50, 115),
        "Table 2.1: Health Plan Coverage Matrix",
        fontsize=11,
        fontname="helv",
        color=(0.1, 0.2, 0.4),
    )
    headers = [
        "Plan Option",
        "Individual Deductible",
        "Family Deductible",
        "Co-Pay (PCP)",
        "Out-of-Pocket Max",
    ]
    x_pos = [55, 140, 250, 360, 460]

    rect_table = fitz.Rect(50, 130, 545, 230)
    page.draw_rect(rect_table, color=(0.7, 0.7, 0.7), fill=(0.96, 0.98, 0.99))

    for x, h in zip(x_pos, headers, strict=False):
        page.insert_text((x, 148), h, fontsize=8, fontname="helv", color=(0.1, 0.2, 0.5))

    rows = [
        ["Standard PPO", "$500 / yr", "$1,000 / yr", "$25.00", "$3,500"],
        ["Premier Plus", "$250 / yr", "$500 / yr", "$15.00", "$2,000"],
        ["HDHP / HSA", "$1,500 / yr", "$3,000 / yr", "0% after ded.", "$4,000"],
        ["Executive Gold", "$0 (Zero)", "$0 (Zero)", "$0.00", "$1,000"],
    ]
    y = 168
    for row in rows:
        page.draw_line(fitz.Point(50, y - 5), fitz.Point(545, y - 5), color=(0.85, 0.85, 0.85))
        for x, val in zip(x_pos, row, strict=False):
            page.insert_text((x, y + 10), val, fontsize=8, fontname="helv")
        y += 18

    # Table 2: 401(k) Matching Schedule
    page.insert_text(
        (50, 260),
        "Table 2.2: 401(k) Retirement Employer Matching Schedule",
        fontsize=11,
        fontname="helv",
        color=(0.1, 0.2, 0.4),
    )
    headers2 = ["Employee Contribution", "Company Match %", "Effective Match", "Vesting Period"]
    x_pos2 = [55, 200, 330, 440]

    rect_table2 = fitz.Rect(50, 275, 545, 380)
    page.draw_rect(rect_table2, color=(0.7, 0.7, 0.7), fill=(0.98, 0.96, 0.96))
    for x, h in zip(x_pos2, headers2, strict=False):
        page.insert_text((x, 293), h, fontsize=8, fontname="helv", color=(0.4, 0.1, 0.1))

    rows2 = [
        ["First 3% of salary", "100% Match", "3.0% of Base", "Immediate (100%)"],
        ["Next 2% of salary (4%-5%)", "50% Match", "1.0% of Base", "Immediate (100%)"],
        ["Contributions > 5%", "0% Match", "0.0% of Base", "N/A"],
        ["Total Max Employer Match", "—", "4.0% of Base", "Immediate (100%)"],
    ]
    y = 313
    for row in rows2:
        page.draw_line(fitz.Point(50, y - 5), fitz.Point(545, y - 5), color=(0.85, 0.85, 0.85))
        for x, val in zip(x_pos2, row, strict=False):
            page.insert_text((x, y + 10), val, fontsize=8, fontname="helv")
        y += 18

    doc.save(str(output_path))
    doc.close()


def generate_visual_charts_pdf(output_path: Path) -> None:
    """Generate a document with embedded vector drawings, figure boxes, and workflow charts."""
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)

    page.insert_text(
        (50, 60),
        "ORGANIZATIONAL ESCALATION & APPROVAL FLOWCHART",
        fontsize=15,
        fontname="helv",
        color=(0.1, 0.2, 0.5),
    )
    page.insert_text(
        (50, 85),
        "Figure 1: Grievance and Dispute Escalation Workflow",
        fontsize=11,
        fontname="helv",
        color=(0.2, 0.2, 0.2),
    )

    # Draw workflow chart boxes
    boxes = [
        (
            "Step 1: Informal Discussion with Direct Manager (Within 5 Days)",
            110,
            (0.9, 0.95, 1.0),
            (0.1, 0.3, 0.7),
        ),
        (
            "Step 2: Formal Written Submission to Department Head (Within 10 Days)",
            180,
            (0.95, 0.95, 1.0),
            (0.2, 0.2, 0.6),
        ),
        (
            "Step 3: People Operations Independent Review Committee (Within 15 Days)",
            250,
            (1.0, 0.95, 0.9),
            (0.7, 0.3, 0.1),
        ),
        (
            "Step 4: Final Executive Binding Determination (Within 30 Days)",
            320,
            (0.9, 1.0, 0.9),
            (0.1, 0.6, 0.2),
        ),
    ]

    for label, y, bg_color, border_color in boxes:
        rect = fitz.Rect(70, y, 525, y + 45)
        page.draw_rect(rect, color=border_color, fill=bg_color, width=1.5)
        page.insert_text((85, y + 27), label, fontsize=9, fontname="helv", color=(0.1, 0.1, 0.1))
        if y < 320:
            # Draw arrow line
            page.draw_line(
                fitz.Point(297, y + 45), fitz.Point(297, y + 70), color=(0.4, 0.4, 0.4), width=1.5
            )

    page.insert_text(
        (50, 400),
        "Figure Caption 1.1: Disputes must exhaust Level 1 and 2 before HR committee escalation.",
        fontsize=9,
        fontname="helv",
        color=(0.4, 0.4, 0.4),
    )
    doc.save(str(output_path))
    doc.close()


def generate_scanned_signed_pdf(output_path: Path) -> None:
    """Generate a simulated scanned PDF with an embedded raster image and signature block."""
    # Create synthetic scanned image using Pillow
    img = Image.new("RGB", (800, 1100), color=(248, 247, 242))
    draw = ImageDraw.Draw(img)

    # Draw document header and text onto image
    draw.text((60, 80), "EMPLOYEE POLICY ACKNOWLEDGEMENT & ATTESTATION", fill=(40, 40, 40))
    draw.text(
        (60, 130),
        "I hereby acknowledge receipt and understanding of the Enterprise Code of Conduct.",
        fill=(60, 60, 60),
    )
    draw.text(
        (60, 160),
        "I certify compliance with all confidentiality, IP assignment, and anti-harassment policies.",
        fill=(60, 60, 60),
    )
    draw.text((60, 220), "Employee Signature: John Doe, Senior Staff Engineer", fill=(10, 30, 120))
    draw.text((60, 250), "Date of Signing: February 14, 2026", fill=(60, 60, 60))
    draw.text(
        (60, 280),
        "Witnessed by HR Representative: Jane Smith, Director of People Ops",
        fill=(60, 60, 60),
    )
    draw.line([(60, 350), (740, 350)], fill=(180, 180, 180), width=1)

    temp_img_path = output_path.parent / "temp_scanned_page.png"
    img.save(temp_img_path)

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_image(fitz.Rect(0, 0, 595, 842), filename=str(temp_img_path))
    doc.save(str(output_path))
    doc.close()
    if temp_img_path.exists():
        temp_img_path.unlink()


def generate_boilerplate_heavy_pdf(output_path: Path, pages_count: int = 15) -> None:
    """Generate a multi-page PDF with repeated headers, footers, and disclaimers on every page."""
    doc = fitz.open()

    for p in range(1, pages_count + 1):
        page = doc.new_page(width=595, height=842)
        # Repeated Header
        page.insert_text(
            (50, 40),
            "CONFIDENTIAL & PROPRIETARY — ACME GLOBAL CORP",
            fontsize=8,
            fontname="helv",
            color=(0.5, 0.5, 0.5),
        )
        page.draw_line(fitz.Point(50, 48), fitz.Point(545, 48), color=(0.8, 0.8, 0.8))

        # Body Content
        page.insert_text(
            (50, 80),
            f"Section {p}: Annual Regulatory Compliance Audit",
            fontsize=13,
            fontname="helv",
            color=(0.1, 0.2, 0.5),
        )
        page.insert_text(
            (50, 105),
            f"This is section {p} of the regulatory adherence report. Sub-article {p}.4 covers data governance.",
            fontsize=10,
            fontname="helv",
        )
        page.insert_text(
            (50, 125),
            f"Audit item {p}-A confirmed zero non-conformances with local labor regulations.",
            fontsize=10,
            fontname="helv",
        )
        page.insert_text(
            (50, 145),
            f"Employee mandatory training completion rate was verified at 99.{p}% for department {p}.",
            fontsize=10,
            fontname="helv",
        )

        # Repeated Footer
        page.draw_line(fitz.Point(50, 800), fitz.Point(545, 800), color=(0.8, 0.8, 0.8))
        page.insert_text(
            (50, 815),
            "Notice: Unauthorized distribution is strictly prohibited under federal law.",
            fontsize=7,
            fontname="helv",
            color=(0.5, 0.5, 0.5),
        )
        page.insert_text(
            (490, 815),
            f"Page {p} of {pages_count}",
            fontsize=8,
            fontname="helv",
            color=(0.3, 0.3, 0.3),
        )

    doc.save(str(output_path))
    doc.close()


def generate_docx_policy(output_path: Path) -> None:
    """Generate a formatted Word DOCX policy document."""
    doc = docx.Document()
    doc.add_heading("Global Executive Travel & Expense Policy", level=0)

    p = doc.add_paragraph("Policy ID: POL-FIN-2026-04 | Effective: Jan 1, 2026")
    p.runs[0].font.color.rgb = RGBColor(100, 100, 100)

    doc.add_heading("1. Purpose and Scope", level=1)
    doc.add_paragraph(
        "This policy defines authorized business travel expenses, per-diem allowances, and reimbursement "
        "procedures for all employees traveling on official company business."
    )

    doc.add_heading("2. Per-Diem Rates & Meal Allowances", level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Travel Tier"
    hdr_cells[1].text = "Daily Meal Cap (USD)"
    hdr_cells[2].text = "Incidentals Allowance"

    data = [
        ("Tier 1: High Cost Metro (NYC, London, Tokyo, SF)", "$120.00 / day", "$25.00 / day"),
        ("Tier 2: Standard Metropolitan Areas", "$85.00 / day", "$15.00 / day"),
        ("Tier 3: Regional / Domestic Non-Metro", "$65.00 / day", "$10.00 / day"),
    ]
    for tier, meal, inc in data:
        row_cells = table.add_row().cells
        row_cells[0].text = tier
        row_cells[1].text = meal
        row_cells[2].text = inc

    doc.add_heading("3. Flight Booking Rules", level=1)
    doc.add_paragraph(
        "Flights under 6 hours duration must be booked in Economy / Standard Main Cabin.",
        style="List Bullet",
    )
    doc.add_paragraph(
        "International flights exceeding 8 continuous hours are eligible for Premium Economy or Business Class.",
        style="List Bullet",
    )

    doc.save(str(output_path))


def generate_xlsx_bonus_matrix(output_path: Path) -> None:
    """Generate a multi-sheet Excel spreadsheet with compensation and bonus formulas."""
    wb = openpyxl.Workbook()
    ws = wb.active
    if ws is None:
        ws = wb.create_sheet("Bonus Structure 2026")
    else:
        ws.title = "Bonus Structure 2026"

    # Headers
    headers = ["Grade", "Role Title", "Target Bonus %", "Max Multiplier", "Min Multiplier"]
    ws.append(headers)

    rows = [
        ["Grade 1", "Associate Specialist", 0.05, 1.5, 0.0],
        ["Grade 2", "Senior Specialist", 0.10, 1.75, 0.0],
        ["Grade 3", "Lead / Staff Engineer", 0.15, 2.0, 0.0],
        ["Grade 4", "Principal / Manager", 0.20, 2.0, 0.0],
        ["Grade 5", "Director / VP", 0.35, 2.5, 0.0],
    ]
    for r in rows:
        ws.append(r)

    # Styling
    header_fill = PatternFill(start_color="1F497D", end_color="1F497D", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    for col in range(1, 6):
        cell = ws.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")

    wb.save(str(output_path))


def generate_pptx_slides(output_path: Path) -> None:
    """Generate a PowerPoint presentation with slide hierarchy."""
    prs = pptx.Presentation()

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    if slide1.shapes.title is not None:
        slide1.shapes.title.text = "Enterprise Diversity, Equity & Inclusion"
    if len(slide1.placeholders) > 1:
        slide1.placeholders[1].text = "2026 Global Strategic Roadmap & Action Plan"

    # Slide 2: Pillars
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    if slide2.shapes.title is not None:
        slide2.shapes.title.text = "Strategic DEI Pillars"
    if len(slide2.placeholders) > 1:
        slide2.placeholders[1].text = (
            "1. Inclusive Talent Acquisition\n"
            "2. Equitable Career Progression\n"
            "3. Global ERG & Community Support"
        )

    prs.save(str(output_path))


def main() -> None:
    CORPUS_DIR.mkdir(parents=True, exist_ok=True)

    print("Generating comprehensive benchmark corpus...")
    generate_comprehensive_hr_handbook(CORPUS_DIR / "01_enterprise_hr_handbook_comprehensive.pdf")
    generate_table_heavy_compensation(CORPUS_DIR / "02_benefits_and_compensation_table_heavy.pdf")
    generate_visual_charts_pdf(CORPUS_DIR / "03_org_hierarchy_and_workflow_charts.pdf")
    generate_scanned_signed_pdf(CORPUS_DIR / "04_scanned_signed_policy_acknowledgement.pdf")
    generate_boilerplate_heavy_pdf(
        CORPUS_DIR / "08_boilerplate_heavy_annual_compliance.pdf", pages_count=10
    )
    generate_docx_policy(CORPUS_DIR / "06_executive_travel_policy_v2.docx")
    generate_xlsx_bonus_matrix(CORPUS_DIR / "05_quarterly_bonus_calculator_matrix.xlsx")
    generate_pptx_slides(CORPUS_DIR / "07_enterprise_diversity_and_inclusion.pptx")
    print(f"Benchmark corpus successfully assembled in {CORPUS_DIR}")


if __name__ == "__main__":
    main()
